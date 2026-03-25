from __future__ import annotations

from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path.cwd()
TEMPLATE = ROOT / "docs" / "App_Description_template_eng_v.1.42.pptx"
APP_NAME = "GoMax Short"
COMPANY_NAME = "GoMax Entertainment Ltd"
OUTPUT = ROOT / f"App_Description_{APP_NAME}_filled.pptx"
SCREENSHOT_DIR = ROOT / "store-screenshots"
UI_DIR = ROOT / "ui-description-screenshots"


def remove_slide(prs: Presentation, index: int) -> None:
    slide_id = prs.slides._sldIdLst[index]
    prs.part.drop_rel(slide_id.rId)
    prs.slides._sldIdLst.remove(slide_id)


def remove_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def clear_and_set_text(shape, lines: Iterable[str], font_size: int = 20, bold: bool = False) -> None:
    text_frame = shape.text_frame
    text_frame.clear()

    for idx, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
        paragraph.text = line
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)


def set_cell_text(cell, text: str, font_size: int = 12, bold: bool = False) -> None:
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold


def find_font(size: int) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        Path("C:/Windows/Fonts/arial.ttf"),
        Path("C:/Windows/Fonts/calibri.ttf"),
    ]
    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), size=size)
    return ImageFont.load_default()


def create_collage(items: list[tuple[Path, str]], output_path: Path, title: str) -> None:
    width, height = 1600, 900
    margin = 40
    gutter = 24
    title_height = 70
    cols, rows = 2, 2
    cell_width = (width - margin * 2 - gutter) // cols
    cell_height = (height - margin * 2 - title_height - gutter) // rows
    image_height = cell_height - 52

    canvas = Image.new("RGB", (width, height), (248, 248, 248))
    draw = ImageDraw.Draw(canvas)
    title_font = find_font(34)
    label_font = find_font(24)

    draw.text((margin, 18), title, fill=(28, 28, 28), font=title_font)

    for idx, (image_path, label) in enumerate(items):
        row = idx // cols
        col = idx % cols
        x = margin + col * (cell_width + gutter)
        y = margin + title_height + row * (cell_height + gutter)

        image = Image.open(image_path).convert("RGB")
        fitted = ImageOps.fit(image, (cell_width, image_height), method=Image.Resampling.LANCZOS)
        canvas.paste(fitted, (x, y))

        label_top = y + image_height
        draw.rectangle((x, label_top, x + cell_width, y + cell_height), fill=(255, 255, 255))
        draw.rectangle((x, y, x + cell_width, y + cell_height), outline=(210, 210, 210), width=2)
        draw.text((x + 16, label_top + 12), label, fill=(40, 40, 40), font=label_font)

    canvas.save(output_path, quality=90)


def add_centered_box(slide, left, top, width, height, text: str, fill_rgb: tuple[int, int, int]) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shape.line.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_arrow_text(slide, left, top, width, height, text: str) -> None:
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)


def main() -> None:
    if not TEMPLATE.exists():
        raise FileNotFoundError(f"Template not found: {TEMPLATE}")

    # Create collages using available screenshots
    usage_collage = UI_DIR / "08-usage-scenario-collage.jpg"
    menu_collage = UI_DIR / "09-menu-function-collage.jpg"
    navigation_collage = UI_DIR / "10-navigation-flow-collage.jpg"

    create_collage(
        [
            (SCREENSHOT_DIR / "01-home-hero.jpg", "1. Home screen"),
            (SCREENSHOT_DIR / "02-home-featured.jpg", "2. Featured content"),
            (SCREENSHOT_DIR / "03-discover-top.jpg", "3. Discover page"),
            (SCREENSHOT_DIR / "04-discover-categories.jpg", "4. Category browsing"),
        ],
        usage_collage,
        f"{APP_NAME} - Usage Scenario",
    )

    create_collage(
        [
            (SCREENSHOT_DIR / "01-home-hero.jpg", "Home"),
            (SCREENSHOT_DIR / "03-discover-top.jpg", "Discover"),
            (UI_DIR / "03-history-page.jpg", "My List"),
            (UI_DIR / "04-settings-page.jpg", "Categories"),
        ],
        menu_collage,
        f"{APP_NAME} - Main Screens",
    )

    create_collage(
        [
            (SCREENSHOT_DIR / "01-home-hero.jpg", "1. Home navigation"),
            (SCREENSHOT_DIR / "03-discover-top.jpg", "2. Discover content"),
            (UI_DIR / "03-history-page.jpg", "3. Watch history"),
            (UI_DIR / "04-settings-page.jpg", "4. Category browsing"),
        ],
        navigation_collage,
        f"{APP_NAME} - Navigation Flow",
    )

    prs = Presentation(str(TEMPLATE))
    remove_slide(prs, 0)

    cover = prs.slides[0]
    cover.shapes[0].text = f"{APP_NAME} Description"
    cover.shapes[1].text = COMPANY_NAME

    revision = prs.slides[1]
    revision.shapes[0].text = "Revision History"
    table = revision.shapes[1].table
    set_cell_text(table.cell(0, 0), "Version", 12, True)
    set_cell_text(table.cell(0, 1), "Date", 12, True)
    set_cell_text(table.cell(0, 2), "Description", 12, True)
    set_cell_text(table.cell(0, 3), "Author", 12, True)
    set_cell_text(table.cell(1, 0), "1.0")
    set_cell_text(table.cell(1, 1), "March 25, 2026")
    set_cell_text(table.cell(1, 2), f"Initial Samsung TV Seller Office UI description for {APP_NAME}.")
    set_cell_text(table.cell(1, 3), "GoMax Team")
    set_cell_text(table.cell(2, 0), "1.0.1")
    set_cell_text(table.cell(2, 1), "March 25, 2026")
    set_cell_text(table.cell(2, 2), "Updated with actual application screenshots and navigation flow.")
    set_cell_text(table.cell(2, 3), "GoMax Team")

    contents = prs.slides[2]
    contents.shapes[0].text = "Contents"
    clear_and_set_text(
        contents.shapes[1],
        [
            "UI Structure",
            "Usage Scenario",
            "Menu & Function Description",
            "Navigation Flow",
            "Key Policy",
            "Language Support",
        ],
        font_size=24,
    )

    ui_overview = prs.slides[3]
    ui_overview.shapes[0].text = "UI Structure"
    clear_and_set_text(
        ui_overview.shapes[1],
        [
            "Whole UI Structure",
            "",
            f"{APP_NAME} is a Samsung TV short drama streaming application",
            "for browsing and watching free short drama series.",
            "Only currently available user flows are described in this document.",
            "",
            "Primary user flow:",
            "Home -> Detail -> Player",
            "Discover -> Detail -> Player",
            "",
            "Main screens in the current release:",
            "- Home: featured banner, popular dramas, and quick access cards",
            "- Discover: category browsing, trending content, and new releases",
            "- Detail: drama information, episode list, and related content",
            "- Player: video playback with episode navigation",
            "- My List: watch history and favorites",
            "- Categories: browse by drama type",
        ],
        font_size=18,
    )

    flow_slide = prs.slides[4]
    flow_slide.shapes[0].text = "UI Structure - Flow Graph"
    remove_shape(flow_slide.shapes[1])
    add_centered_box(flow_slide, Inches(1.9), Inches(1.9), Inches(1.7), Inches(0.65), "Home", (0x3F, 0xA7, 0x7A))
    add_arrow_text(flow_slide, Inches(3.75), Inches(1.95), Inches(0.5), Inches(0.5), "<->")
    add_centered_box(flow_slide, Inches(4.3), Inches(1.9), Inches(1.95), Inches(0.65), "Discover", (0x3F, 0xA7, 0x7A))
    add_arrow_text(flow_slide, Inches(3.95), Inches(2.75), Inches(0.3), Inches(0.5), "v")
    add_centered_box(flow_slide, Inches(3.25), Inches(3.2), Inches(1.8), Inches(0.7), "Detail", (0xE0, 0x8A, 0x3B))
    add_arrow_text(flow_slide, Inches(4.0), Inches(3.95), Inches(0.3), Inches(0.5), "v")
    add_centered_box(flow_slide, Inches(3.05), Inches(4.35), Inches(2.2), Inches(0.75), "Player", (0xB9, 0x4E, 0x5D))

    depth_slide = prs.slides[5]
    depth_slide.shapes[0].text = "UI Structure - Depth Navigation"
    remove_shape(depth_slide.shapes[1])
    textbox = depth_slide.shapes.add_textbox(Inches(1.15), Inches(2.0), Inches(7.1), Inches(4.1))
    clear_and_set_text(
        textbox,
        [
            "Depth 1 : Home, Discover",
            "Depth 2 : Selected drama detail page",
            "Depth 3 : Episode playback page",
            "",
            "Global UI elements:",
            "- Header navigation bar",
            "- Focus highlight for remote-control navigation",
            "",
            "Contextual UI elements:",
            "- Episode list on Detail page",
            "- Video area and episode list on Player page",
        ],
        font_size=20,
    )

    usage_summary = prs.slides[6]
    usage_summary.shapes[0].text = "Usage Scenario"
    clear_and_set_text(
        usage_summary.shapes[1],
        [
            "Use Case 1: Browse and watch short dramas.",
            "",
            "1. Navigate between Home, Discover, My List, and Categories with directional keys.",
            "2. Browse featured dramas on Home or explore categories in Discover.",
            "3. Select a drama card with the Enter key to view details.",
            "4. Review drama information, synopsis, and episode list on Detail screen.",
            "5. Select Play button or choose a specific episode to start watching.",
            "6. Control playback with remote: play/pause, seek, next/previous episode.",
            "7. Press Return to navigate back or access episode list.",
            "8. Add dramas to favorites or view watch history in My List.",
            "",
            "No account sign-in, activation, or purchase is required - all content is free!",
        ],
        font_size=18,
    )

    usage_sample = prs.slides[7]
    usage_sample.shapes[0].text = "Usage Scenario"
    clear_and_set_text(
        usage_sample.shapes[1],
        [
            "Use Case Title: Navigate through main application screens",
            "Path 1: Home -> Browse featured content",
            "Path 2: Discover -> Explore categories and trending content",
            "Path 3: My List -> View watch history and favorites",
            "Path 4: Categories -> Browse dramas by genre",
            "",
            "Supporting notes:",
            "- The app follows Samsung TV navigation patterns.",
            "- All content is completely free to watch.",
            "- No account registration or payment required.",
            "",
            "No account login is required.",
            "No in-app purchase is required.",
        ],
        font_size=18,
    )
    usage_picture = usage_sample.shapes[2]
    left, top, width, height = usage_picture.left, usage_picture.top, usage_picture.width, usage_picture.height
    remove_shape(usage_picture)
    usage_sample.shapes.add_picture(str(navigation_collage), left, top, width=width, height=height)

    menu_slide = prs.slides[8]
    menu_slide.shapes[0].text = "Menu & Function Description"
    clear_and_set_text(
        menu_slide.shapes[1],
        [
            "Screen descriptions in the current release:",
            "",
            "Home: features hero banner, featured dramas, and quick access cards.",
            "Discover: displays category browsing, trending content, and new releases.",
            "Detail: shows drama poster, synopsis, episode grid, and related content.",
            "Player: full-screen video with episode sidebar and playback controls.",
            "My List: watch history and favorite dramas collection.",
            "Categories: browse dramas by genre and type.",
            "",
            "All content is completely free - no subscriptions or payments required!",
        ],
        font_size=18,
    )
    menu_picture = menu_slide.shapes[2]
    left, top, width, height = menu_picture.left, menu_picture.top, menu_picture.width, menu_picture.height
    remove_shape(menu_picture)
    menu_slide.shapes.add_picture(str(menu_collage), left, top, width=width, height=height)

    key_slide = prs.slides[9]
    key_slide.shapes[0].text = "Key Policy"
    clear_and_set_text(
        key_slide.shapes[1],
        [
            "The current release follows Samsung TV remote control patterns.",
            "Home page: Navigate between main sections and featured content.",
            "Discover / Detail / Player pages: Return navigates to the previous screen.",
            "Player controls: Dedicated media keys control playback functions.",
            "Navigation: Directional keys move focus between interactive elements.",
        ],
        font_size=18,
    )
    key_table = key_slide.shapes[2].table
    rows = [
        ("Button", "Action", "Remarks"),
        ("ENTER", "Select the focused card, button, or episode", "Standard selection behavior"),
        ("UP / DOWN", "Move focus vertically", "Navigate between sections"),
        ("LEFT / RIGHT", "Move focus horizontally", "Navigate within rows"),
        ("RETURN", "Go to the previous screen", "Standard back navigation"),
        ("PLAY / PAUSE", "Control video playback", "Media control functions"),
        ("FF / RW", "Seek forward/backward 10 seconds", "Quick navigation in video"),
        ("EXIT", "Close the application", "Platform default behavior"),
        ("COLOR / CHANNEL / NUMBER KEYS", "N/R", "No custom mapping in current release"),
    ]
    for r_idx, row in enumerate(rows):
        for c_idx, value in enumerate(row):
            set_cell_text(key_table.cell(r_idx, c_idx), value, 11, r_idx == 0)

    lang_slide = prs.slides[10]
    lang_slide.shapes[0].text = "How to Change Languages"
    lang_table = lang_slide.shapes[1].table
    set_cell_text(lang_table.cell(0, 0), "Items", 12, True)
    set_cell_text(lang_table.cell(0, 1), "Contents", 12, True)
    set_cell_text(lang_table.cell(1, 0), "Language support")
    set_cell_text(
        lang_table.cell(1, 1),
        "This version supports English only. No in-app language change menu is provided.",
    )

    prs.save(str(OUTPUT))
    print(f"PowerPoint presentation saved: {OUTPUT}")


if __name__ == "__main__":
    main()